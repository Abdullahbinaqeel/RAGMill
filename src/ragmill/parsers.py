"""
Format-specific text extractors.

Each extractor takes a file path and returns its plain-text content. All
third-party imports are lazy so the core package stays installable with zero
dependencies — a format's extractor only needs its extra when that format is
actually encountered.

Extras:
  - pdf    → pypdf                          (.pdf, digital text layer)
  - docx   → python-docx                    (.docx)
  - office → beautifulsoup4, striprtf,      (.html/.htm, .rtf, .xlsx, .pptx)
             openpyxl, python-pptx
  - ocr    → pytesseract, pillow            (.png/.jpg/... and scanned .pdf)
             + system binaries: `tesseract`, and `pdftoppm` (poppler) for PDFs

OCR is English-only by default (lang="eng").
"""

import logging
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

DEFAULT_OCR_LANG = "eng"
# A PDF whose extracted text layer is shorter than this is treated as
# scanned/image-only and routed to OCR (when enabled).
PDF_OCR_MIN_CHARS = 10

# How to install the system binaries OCR needs, per platform. pip cannot supply
# these, so the message has to name a real command for the user's OS — a
# brew-only hint is useless to the Windows and Linux majority.
_BINARY_HINTS = {
    "tesseract": {
        "win32": "download the installer from https://github.com/UB-Mannheim/tesseract/wiki "
        "and make sure its folder is on PATH",
        "darwin": "brew install tesseract",
        "linux": "apt-get install tesseract-ocr",
    },
    "pdftoppm": {
        "win32": "download poppler from https://github.com/oschwartz10612/poppler-windows/releases "
        "and add its bin/ folder to PATH",
        "darwin": "brew install poppler",
        "linux": "apt-get install poppler-utils",
    },
}


# Byte-order marks, longest first — UTF-32 LE's BOM starts with UTF-16 LE's, so
# order matters. The codecs here are the BOM-consuming ones on purpose: bare
# "utf-16-le" leaves U+FEFF at the start of the text, which then leads the first
# chunk. "utf-16"/"utf-32" read the BOM for endianness and strip it; "utf-8-sig"
# strips the UTF-8 BOM.
_BOMS = (
    (b"\x00\x00\xfe\xff", "utf-32"),
    (b"\xff\xfe\x00\x00", "utf-32"),
    (b"\xef\xbb\xbf", "utf-8-sig"),
    (b"\xfe\xff", "utf-16"),
    (b"\xff\xfe", "utf-16"),
)

# Tried in order for BOM-less files that are not valid UTF-8. cp1252 is what
# Windows Notepad's "ANSI" produces; latin-1 cannot fail, so it terminates the
# chain rather than letting a file go unread.
_TEXT_FALLBACKS = ("cp1252", "latin-1")


def read_text_file(path: str) -> str:
    """Read a text file, detecting its encoding instead of assuming UTF-8.

    Reading everything as UTF-8 with errors="ignore" corrupts files silently,
    which matters most on Windows because that is what Notepad emits:

      - "Unicode" (UTF-16) decoded as UTF-8 yields NUL-interleaved mojibake
        ("T\\x00h\\x00e\\x00"), so chunks embed as noise and never match a query.
      - "ANSI" (cp1252) loses every non-ASCII byte outright, turning "costs £50"
        into "costs 50" — a silent change of meaning, not a visible failure.

    Order: honour a BOM if present, else strict UTF-8, else the fallbacks. A
    fallback is logged, because guessing an encoding is worth telling the user.
    """
    with open(path, "rb") as f:
        raw = f.read()

    for bom, codec in _BOMS:
        if raw.startswith(bom):
            return raw.decode(codec)

    try:
        text = raw.decode("utf-8")
    except UnicodeDecodeError:
        pass
    else:
        # BOM-less UTF-16 can decode as UTF-8 (NUL is valid U+0000) and produce
        # exactly the mojibake above. Interior NULs never occur in real text.
        if "\x00" not in text:
            return text
        for codec in ("utf-16-le", "utf-16-be"):
            try:
                candidate = raw.decode(codec)
            except UnicodeDecodeError:
                continue
            if "\x00" not in candidate:
                logger.warning(
                    "%s: no BOM but content looks like %s; decoded as %s.",
                    os.path.basename(path),
                    codec,
                    codec,
                )
                return candidate
        return text

    for codec in _TEXT_FALLBACKS:
        try:
            text = raw.decode(codec)
        except UnicodeDecodeError:
            continue
        logger.warning(
            "%s: not valid UTF-8, decoded as %s. Some characters may be wrong — "
            "re-save the file as UTF-8 if the text looks incorrect.",
            os.path.basename(path),
            codec,
        )
        return text

    raise UnicodeDecodeError(  # pragma: no cover — latin-1 above cannot fail
        "utf-8", raw, 0, len(raw), f"could not decode {path} as UTF-8, cp1252, or latin-1"
    )


def _install_hint(binary: str) -> str:
    """Platform-appropriate instructions for installing a system binary."""
    hints = _BINARY_HINTS[binary]
    if sys.platform.startswith("win"):
        return hints["win32"]
    if sys.platform == "darwin":
        return hints["darwin"]
    return hints["linux"]


# ── PDF (digital text layer, with OCR fallback for scans) ────────────────────


def extract_pdf_text(
    path: str,
    *,
    enable_ocr: bool = True,
    ocr_lang: str = DEFAULT_OCR_LANG,
    min_chars: int = PDF_OCR_MIN_CHARS,
) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ImportError(
            "PDF support requires the 'pdf' extra. Install it with: pip install ragmill[pdf]"
        ) from exc

    reader = PdfReader(path)
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(page.strip() for page in pages if page.strip())

    if len(text) >= min_chars or not enable_ocr:
        return text

    # No usable text layer → scanned/image PDF. Rasterise then OCR.
    return _ocr_pdf(path, ocr_lang)


def _ocr_pdf(path: str, ocr_lang: str) -> str:
    if not shutil.which("pdftoppm"):
        raise RuntimeError(
            "This PDF has no text layer, so it looks scanned and needs OCR — which "
            "requires poppler's `pdftoppm` to rasterise the pages. To install it, "
            f"{_install_hint('pdftoppm')}. To skip OCR and accept empty text for "
            "scanned PDFs instead, call extract_pdf_text(..., enable_ocr=False)."
        )
    logger.info("%s: no text layer, running OCR (%s)…", os.path.basename(path), ocr_lang)
    with tempfile.TemporaryDirectory() as tmp:
        prefix = os.path.join(tmp, "page")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "300", str(path), prefix],
            check=True,
            capture_output=True,
        )
        parts = [_ocr_image_file(png, ocr_lang) for png in sorted(Path(tmp).glob("page*.png"))]
    return "\n\n".join(p for p in parts if p).strip()


# ── DOCX ──────────────────────────────────────────────────────────────────────


def extract_docx_text(path: str) -> str:
    try:
        import docx
    except ImportError as exc:
        raise ImportError(
            "DOCX support requires the 'docx' extra. Install it with: pip install ragmill[docx]"
        ) from exc

    document = docx.Document(path)
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


# ── Images + shared OCR helper ────────────────────────────────────────────────


def _ocr_image_file(image_path, ocr_lang: str) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise ImportError(
            'OCR support requires the "ocr" extra: pip install "ragmill[ocr]". '
            f"It also needs the system `tesseract` binary — to install that, {_install_hint('tesseract')}."
        ) from exc
    if not shutil.which("tesseract"):
        raise RuntimeError(
            "The `tesseract` OCR binary is not on PATH. To install it, "
            f"{_install_hint('tesseract')}."
        )
    with Image.open(image_path) as img:
        return pytesseract.image_to_string(img, lang=ocr_lang).strip()


def extract_image_text(path: str, *, ocr_lang: str = DEFAULT_OCR_LANG) -> str:
    return _ocr_image_file(path, ocr_lang)


# ── HTML ──────────────────────────────────────────────────────────────────────


def extract_html_text(path: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise ImportError(
            "HTML support requires the 'office' extra. Install it with: pip install ragmill[office]"
        ) from exc
    soup = BeautifulSoup(read_text_file(path), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n").strip()


# ── RTF ───────────────────────────────────────────────────────────────────────


def extract_rtf_text(path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as exc:
        raise ImportError(
            "RTF support requires the 'office' extra. Install it with: pip install ragmill[office]"
        ) from exc
    return rtf_to_text(read_text_file(path)).strip()


# ── XLSX ──────────────────────────────────────────────────────────────────────


def extract_xlsx_text(path: str) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ImportError(
            "XLSX support requires the 'office' extra. Install it with: pip install ragmill[office]"
        ) from exc

    workbook = load_workbook(filename=path, read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    workbook.close()
    return "\n".join(parts).strip()


# ── PPTX ──────────────────────────────────────────────────────────────────────


def extract_pptx_text(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError(
            "PPTX support requires the 'office' extra. Install it with: pip install ragmill[office]"
        ) from exc

    prs = Presentation(path)
    parts = []
    for index, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts).strip()


# ── CSV / TSV (no extra needed — stdlib) ─────────────────────────────────────


def extract_csv_text(path: str) -> str:
    import csv

    dialect = "excel-tab" if os.path.splitext(path)[1].lower() == ".tsv" else "excel"
    rows = []
    # splitlines() keeps csv's newline handling intact after decoding.
    for row in csv.reader(read_text_file(path).splitlines(), dialect=dialect):
        cells = [c.strip() for c in row if c.strip()]
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows).strip()
