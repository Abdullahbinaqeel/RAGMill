"""
Comprehensive integration test: every supported file type through the full
pipeline (parse → chunk → embed → store).

This ensures no format silently breaks when its parser or a dependency
changes.  OCR tests are conditional on the tesseract binary.
"""

import shutil

import pytest

from ragmill import RAGEngine

HAS_TESSERACT = shutil.which("tesseract") is not None

# Text layers used by the PDF fixtures below. Each must stay longer than
# parsers.PDF_OCR_MIN_CHARS — see test_pdf_fixtures_clear_the_ocr_threshold.
PDF_FIXTURE_TEXTS = ("pdf content", "pdf body content", "zeta eta theta")

# ── Helpers ───────────────────────────────────────────────────────────────────


def _font():
    from PIL import ImageFont

    for p in [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]:
        try:
            return ImageFont.truetype(p, 36)
        except OSError:
            continue
    return ImageFont.load_default()


# ── Individual format tests ───────────────────────────────────────────────────


def test_txt(tmp_path):
    p = tmp_path / "a.txt"
    p.write_text("hello from txt")
    engine = RAGEngine()
    docs = list(engine.stream_directory(str(tmp_path)))
    assert len(docs) == 1
    assert "hello from txt" in docs[0]["raw_content"]


def test_md(tmp_path):
    p = tmp_path / "b.md"
    p.write_text("# Title\n\nmarkdown body")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "markdown body" in docs["b.md"]


def test_log(tmp_path):
    p = tmp_path / "app.log"
    p.write_text("INFO: system started\nERROR: disk full")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "disk full" in docs["app.log"]


def test_rst(tmp_path):
    p = tmp_path / "doc.rst"
    p.write_text("=====\nTitle\n=====\n\nrst body")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "rst body" in docs["doc.rst"]


def test_csv(tmp_path):
    p = tmp_path / "data.csv"
    p.write_text("col1,col2\nval1,val2\n")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "val1 | val2" in docs["data.csv"]


def test_tsv(tmp_path):
    p = tmp_path / "data.tsv"
    p.write_text("col1\tcol2\nval1\tval2\n")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "val1 | val2" in docs["data.tsv"]


def test_html(tmp_path):
    pytest.importorskip("bs4")
    p = tmp_path / "page.html"
    p.write_text("<html><body><p>html content</p></body></html>")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "html content" in docs["page.html"]


def test_htm(tmp_path):
    pytest.importorskip("bs4")
    p = tmp_path / "page.htm"
    p.write_text("<html><body><p>htm content</p></body></html>")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "htm content" in docs["page.htm"]


def test_rtf(tmp_path):
    pytest.importorskip("striprtf")
    p = tmp_path / "note.rtf"
    p.write_text(r"{\rtf1\ansi rtf content\par}")
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "rtf content" in docs["note.rtf"]


def test_xlsx(tmp_path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    p = tmp_path / "sheet.xlsx"
    wb = Workbook()
    wb.active.append(["name", "value"])
    wb.save(str(p))
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "name | value" in docs["sheet.xlsx"]


def test_pptx(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    p = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "slide content"
    prs.save(str(p))
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "slide content" in docs["deck.pptx"]


def test_docx(tmp_path):
    pytest.importorskip("docx")
    import docx

    p = tmp_path / "letter.docx"
    doc = docx.Document()
    doc.add_paragraph("docx content")
    doc.save(str(p))
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "docx content" in docs["letter.docx"]


def test_pdf(tmp_path):
    pytest.importorskip("pypdf")
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas

    p = tmp_path / "doc.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "pdf content")
    c.save()
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "pdf content" in docs["doc.pdf"]


def test_pdf_fixtures_clear_the_ocr_threshold():
    """Every PDF fixture here must have a text layer longer than PDF_OCR_MIN_CHARS.

    Below that, extract_pdf_text() concludes the PDF is scanned and falls back to
    OCR, which needs poppler's `pdftoppm` and tesseract on the host. Fixtures
    under the threshold therefore pass on a dev machine with brew-installed
    binaries and fail in CI, which has neither.
    """
    from ragmill.parsers import PDF_OCR_MIN_CHARS

    for text in PDF_FIXTURE_TEXTS:
        assert len(text) > PDF_OCR_MIN_CHARS, (
            f"PDF fixture {text!r} is {len(text)} chars, at or under the "
            f"{PDF_OCR_MIN_CHARS}-char OCR threshold; it would exercise the OCR "
            "fallback instead of the text-layer path. Lengthen the fixture."
        )


@pytest.mark.skipif(not HAS_TESSERACT, reason="requires tesseract binary")
def test_png(tmp_path):
    pytest.importorskip("pytesseract")
    pytest.importorskip("PIL")
    from PIL import Image, ImageDraw

    p = tmp_path / "scan.png"
    img = Image.new("RGB", (400, 100), "white")
    ImageDraw.Draw(img).text((20, 30), "ocr text", fill="black", font=_font())
    img.save(str(p))
    engine = RAGEngine()
    docs = {d["filename"]: d["raw_content"] for d in engine.stream_directory(str(tmp_path))}
    assert "ocr text" in docs["scan.png"].lower()


# ── Full pipeline: all types → embeddings ─────────────────────────────────────


@pytest.mark.parametrize(
    "filename, content_factory, needs_extras",
    [
        ("a.txt", lambda p: p.write_text("hello"), []),
        ("b.md", lambda p: p.write_text("# M\n\nbody"), []),
        ("c.log", lambda p: p.write_text("line1"), []),
        ("d.rst", lambda p: p.write_text("====\nT\n====\n\nbody"), []),
        ("e.csv", lambda p: p.write_text("x,y\n1,2\n"), []),
        ("f.tsv", lambda p: p.write_text("x\ty\n1\t2\n"), []),
        ("g.html", lambda p: p.write_text("<p>html</p>"), ["bs4"]),
        ("h.htm", lambda p: p.write_text("<p>htm</p>"), ["bs4"]),
        ("i.rtf", lambda p: p.write_text(r"{\rtf1 rtf}"), ["striprtf"]),
    ],
)
def test_plaintext_types_yield_chunks(tmp_path, filename, content_factory, needs_extras):
    for mod in needs_extras:
        pytest.importorskip(mod)
    content_factory(tmp_path / filename)
    engine = RAGEngine(chunk_size=500, overlap=50)
    payloads = engine.execute_pipeline(str(tmp_path))
    names = {p["metadata"]["filename"] for p in payloads}
    assert filename in names, f"{filename} did not produce any chunks"


@pytest.mark.parametrize(
    "filename, needs_extras",
    [
        ("sheet.xlsx", ["openpyxl"]),
        ("deck.pptx", ["pptx"]),
        ("letter.docx", ["docx"]),
        ("doc.pdf", ["pypdf", "reportlab"]),
    ],
)
def test_binary_types_yield_chunks(tmp_path, filename, needs_extras):
    for mod in needs_extras:
        pytest.importorskip(mod)
    if filename.endswith(".xlsx"):
        from openpyxl import Workbook

        wb = Workbook()
        wb.active.append(["a", "b"])
        wb.save(str(tmp_path / filename))
    elif filename.endswith(".pptx"):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "content"
        prs.save(str(tmp_path / filename))
    elif filename.endswith(".docx"):
        import docx

        doc = docx.Document()
        doc.add_paragraph("content")
        doc.save(str(tmp_path / filename))
    elif filename.endswith(".pdf"):
        from reportlab.pdfgen import canvas

        c = canvas.Canvas(str(tmp_path / filename))
        # Must exceed PDF_OCR_MIN_CHARS, or extract_pdf_text treats the PDF as
        # scanned and falls back to OCR — which needs poppler and tesseract, so
        # this test would then fail on any machine without them.
        c.drawString(100, 750, "pdf body content")
        c.save()
    engine = RAGEngine(chunk_size=500, overlap=50)
    payloads = engine.execute_pipeline(str(tmp_path))
    names = {p["metadata"]["filename"] for p in payloads}
    assert filename in names, f"{filename} did not produce any chunks"


def test_unsupported_extension_skipped(tmp_path):
    (tmp_path / "ignore.bin").write_bytes(b"\x00\x01")
    engine = RAGEngine()
    docs = list(engine.stream_directory(str(tmp_path)))
    assert len(docs) == 0


def test_empty_file_skipped_with_warning(tmp_path, caplog):
    import logging

    (tmp_path / "empty.txt").write_text("   \n\n")
    engine = RAGEngine()
    with caplog.at_level(logging.WARNING):
        docs = list(engine.stream_directory(str(tmp_path)))
    assert len(docs) == 0
    assert any("No extractable text" in r.message for r in caplog.records)


def test_mixed_directory_processes_all_formats(tmp_path):
    """A directory with one file of each type; every format should produce a chunk."""
    pytest.importorskip("bs4")
    pytest.importorskip("openpyxl")
    pytest.importorskip("pptx")
    pytest.importorskip("docx")
    pytest.importorskip("pypdf")
    pytest.importorskip("reportlab")
    from openpyxl import Workbook
    from pptx import Presentation
    import docx
    from reportlab.pdfgen import canvas

    (tmp_path / "a.txt").write_text("alpha")
    (tmp_path / "b.md").write_text("# M\nbeta")
    (tmp_path / "c.csv").write_text("x,y\n1,2\n")
    (tmp_path / "d.html").write_text("<p>gamma</p>")
    (tmp_path / "e.xlsx").write_bytes(b"")
    wb = Workbook()
    wb.active.append(["v"])
    wb.save(str(tmp_path / "e.xlsx"))
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "delta"
    prs.save(str(tmp_path / "f.pptx"))
    doc = docx.Document()
    doc.add_paragraph("epsilon")
    doc.save(str(tmp_path / "g.docx"))
    c = canvas.Canvas(str(tmp_path / "h.pdf"))
    # Longer than PDF_OCR_MIN_CHARS so this exercises the text-layer path rather
    # than the OCR fallback (see test_binary_types_yield_chunks).
    c.drawString(100, 750, "zeta eta theta")
    c.save()
    (tmp_path / "ignore.bin").write_bytes(b"\x00")

    engine = RAGEngine(chunk_size=500, overlap=50)
    payloads = engine.execute_pipeline(str(tmp_path))
    filenames = {p["metadata"]["filename"] for p in payloads}

    assert "a.txt" in filenames
    assert "b.md" in filenames
    assert "c.csv" in filenames
    assert "d.html" in filenames
    assert "e.xlsx" in filenames
    assert "f.pptx" in filenames
    assert "g.docx" in filenames
    assert "h.pdf" in filenames
    assert "ignore.bin" not in filenames


@pytest.mark.skipif(not HAS_TESSERACT, reason="requires tesseract binary")
def test_mixed_includes_image(tmp_path):
    pytest.importorskip("pytesseract")
    pytest.importorskip("PIL")
    from PIL import Image, ImageDraw

    p = tmp_path / "img.png"
    img = Image.new("RGB", (400, 100), "white")
    ImageDraw.Draw(img).text((20, 30), "picture", fill="black", font=_font())
    img.save(str(p))
    (tmp_path / "a.txt").write_text("text")

    engine = RAGEngine()
    payloads = engine.execute_pipeline(str(tmp_path))
    filenames = {p["metadata"]["filename"] for p in payloads}
    assert "img.png" in filenames
    assert "a.txt" in filenames
