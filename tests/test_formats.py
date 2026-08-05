"""
Extraction tests for the extended format support (CSV, HTML, RTF, XLSX, PPTX)
and OCR (images + scanned PDFs).

Optional third-party libs use pytest.importorskip; OCR tests additionally skip
when the `tesseract` / `pdftoppm` binaries aren't installed, so the suite stays
green on minimal environments.
"""

import shutil

import pytest

from ragmill import RAGEngine
from ragmill import parsers

HAS_TESSERACT = shutil.which("tesseract") is not None
HAS_POPPLER = shutil.which("pdftoppm") is not None
FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
]


def _load_font(size):
    from PIL import ImageFont

    for path in FONT_CANDIDATES:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _text_image(text):
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (1000, 260), "white")
    ImageDraw.Draw(img).text((40, 90), text, fill="black", font=_load_font(44))
    return img


def _build_text_pdf(path, body):
    """Minimal single-page PDF with a real text layer (no reportlab needed)."""
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
    ]
    stream = f"BT /F1 14 Tf 72 720 Td ({body}) Tj ET".encode("latin-1")
    objs.append(b"<< /Length %d >>\nstream\n" % len(stream) + stream + b"\nendstream")
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body_bytes in enumerate(objs, 1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % i + body_bytes + b"\nendobj\n"
    xref = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objs) + 1)
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF" % (
        len(objs) + 1,
        xref,
    )
    path.write_bytes(out)


# ── Plain formats ─────────────────────────────────────────────────────────────


def test_extract_csv_text(tmp_path):
    p = tmp_path / "cases.csv"
    p.write_text("case,court\nappeal 35,supreme court\n", encoding="utf-8")
    text = parsers.extract_csv_text(str(p))
    assert "appeal 35 | supreme court" in text


def test_extract_html_text_strips_script_and_style(tmp_path):
    pytest.importorskip("bs4")
    p = tmp_path / "page.html"
    p.write_text(
        "<html><head><style>.a{}</style></head><body><h1>fundamental rights</h1>"
        "<script>var x=1;</script><p>body text</p></body></html>",
        encoding="utf-8",
    )
    text = parsers.extract_html_text(str(p))
    assert "fundamental rights" in text
    assert "body text" in text
    assert "var x" not in text


def test_extract_rtf_text(tmp_path):
    pytest.importorskip("striprtf")
    p = tmp_path / "note.rtf"
    p.write_text(r"{\rtf1\ansi rent restriction ordinance\par}", encoding="utf-8")
    assert "rent restriction ordinance" in parsers.extract_rtf_text(str(p))


def test_extract_xlsx_text(tmp_path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    p = tmp_path / "list.xlsx"
    wb = Workbook()
    wb.active.append(["cause", "list", "entry"])
    wb.save(str(p))
    text = parsers.extract_xlsx_text(str(p))
    assert "cause | list | entry" in text


def test_extract_pptx_text(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    p = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "legal aid overview"
    prs.save(str(p))
    assert "legal aid overview" in parsers.extract_pptx_text(str(p))


# ── OCR ─────────────────────────────────────────────────────────────────────


@pytest.mark.skipif(not HAS_TESSERACT, reason="requires the tesseract binary")
def test_ocr_image(tmp_path):
    pytest.importorskip("pytesseract")
    pytest.importorskip("PIL")
    p = tmp_path / "scan.png"
    _text_image("notice of eviction served").save(str(p))
    text = parsers.extract_image_text(str(p))
    assert "notice of eviction served" in text.lower()


@pytest.mark.skipif(
    not (HAS_TESSERACT and HAS_POPPLER),
    reason="requires tesseract + poppler (pdftoppm)",
)
def test_ocr_scanned_pdf(tmp_path):
    pytest.importorskip("pytesseract")
    pytest.importorskip("pypdf")
    p = tmp_path / "scanned.pdf"
    _text_image("scanned court order text").save(str(p), "PDF", resolution=300)
    text = parsers.extract_pdf_text(str(p))
    assert "scanned court order text" in text.lower()


def test_digital_pdf_uses_text_layer_not_ocr(tmp_path):
    pytest.importorskip("pypdf")
    p = tmp_path / "digital.pdf"
    _build_text_pdf(p, "digital judgment paragraph one")
    # enable_ocr defaults True but the text layer should satisfy min_chars first.
    assert "digital judgment paragraph one" in parsers.extract_pdf_text(str(p))


# ── Engine integration ────────────────────────────────────────────────────────


def test_engine_streams_new_formats(tmp_path):
    pytest.importorskip("bs4")
    pytest.importorskip("openpyxl")
    (tmp_path / "a.csv").write_text("x,y\nfoo,bar\n", encoding="utf-8")
    (tmp_path / "b.html").write_text("<html><body>hello world</body></html>", encoding="utf-8")
    (tmp_path / "ignore.xyz").write_text("nope", encoding="utf-8")

    engine = RAGEngine()
    results = {r["filename"]: r["raw_content"] for r in engine.stream_directory(str(tmp_path))}

    assert "foo | bar" in results["a.csv"]
    assert "hello world" in results["b.html"]
    assert "ignore.xyz" not in results


def test_engine_skips_textless_file_with_warning(tmp_path, caplog):
    import logging

    # An empty supported file deterministically yields no text, exercising the
    # "no extractable text → skip with warning" branch without depending on any
    # OCR binaries (a blank PDF's behaviour varies with tesseract/poppler).
    (tmp_path / "empty.txt").write_text("   \n\n", encoding="utf-8")

    engine = RAGEngine()
    with caplog.at_level(logging.WARNING):
        filenames = [r["filename"] for r in engine.stream_directory(str(tmp_path))]

    assert "empty.txt" not in filenames
    assert any("No extractable text" in rec.message for rec in caplog.records)


# ── System-binary errors name a command for the user's actual OS ─────────────


class TestBinaryInstallHints:
    """`brew install ...` is useless to the Windows and Linux majority.

    pip cannot supply tesseract or pdftoppm, so these messages are the only
    guidance a user gets. They have to name a command that exists on the OS the
    error was raised on.
    """

    @pytest.mark.parametrize("binary", ["tesseract", "pdftoppm"])
    @pytest.mark.parametrize(
        "platform, expected",
        [("win32", "PATH"), ("darwin", "brew install"), ("linux", "apt-get install")],
    )
    def test_hint_matches_platform(self, monkeypatch, binary, platform, expected):
        monkeypatch.setattr(parsers.sys, "platform", platform)
        hint = parsers._install_hint(binary)
        assert expected in hint, f"{binary} on {platform}: {hint!r}"
        if platform != "darwin":
            assert "brew" not in hint, f"{binary} on {platform} still says brew: {hint!r}"

    def test_windows_hints_link_to_a_real_download(self, monkeypatch):
        monkeypatch.setattr(parsers.sys, "platform", "win32")
        assert "UB-Mannheim" in parsers._install_hint("tesseract")
        assert "poppler-windows" in parsers._install_hint("pdftoppm")

    def test_missing_poppler_error_is_windows_aware(self, monkeypatch, tmp_path):
        """The exact failure the CI run hit, but reported to a Windows user."""
        monkeypatch.setattr(parsers.sys, "platform", "win32")
        monkeypatch.setattr(parsers.shutil, "which", lambda _: None)

        with pytest.raises(RuntimeError) as excinfo:
            parsers._ocr_pdf(str(tmp_path / "scan.pdf"), "eng")

        msg = str(excinfo.value)
        assert "poppler-windows" in msg, msg
        assert "brew" not in msg, msg
        # tells the user how to opt out of OCR rather than only how to enable it
        assert "enable_ocr=False" in msg, msg
